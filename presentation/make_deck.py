"""Собирает sparks.pptx — тот же дек, что и sparks.html, но для проектора без браузера.

Запуск:  pip install python-pptx  &&  python make_deck.py
Рядом должен лежать qr-keep-sparks.png.

Пропорции взяты из HTML-версии: там 1rem = 1% высоты слайда, здесь высота
слайда 7.5" = 540pt, поэтому 1rem = 5.4pt. Все размеры пересчитаны по этому
курсу — дек выглядит так же, как в браузере.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x14, 0x14, 0x16)
SURFACE = RGBColor(0x1B, 0x1B, 0x1E)
BORDER = RGBColor(0x2A, 0x2A, 0x2F)
TEXT = RGBColor(0xEC, 0xEC, 0xEE)
MUTED = RGBColor(0x8B, 0x8B, 0x94)
ACCENT = RGBColor(0x8B, 0x7C, 0xF6)

FONT = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
PAD_X = Inches(0.68)  # 9% высоты в пересчёте на ширину — как отступы в HTML
PAD_Y = Inches(0.6)
CONTENT_W = W - 2 * PAD_X

REM = 5.4  # 1% высоты слайда в пунктах


def pt(rem: float) -> Pt:
    """Размер шрифта из HTML-единиц (1rem = 1% высоты слайда)."""
    return Pt(rem * REM)


prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def text(
    s,
    body,
    top,
    size,
    color=TEXT,
    bold=False,
    left=PAD_X,
    width=CONTENT_W,
    align=PP_ALIGN.LEFT,
    spacing=1.2,
):
    """Текстовый блок. Возвращает низ блока, чтобы складывать слайд сверху вниз."""
    height = Pt(size.pt * spacing * (body.count("\n") + 1))
    box = s.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT
    return top + height


def title(s, body):
    return text(s, body, PAD_Y, pt(7.5), bold=True, spacing=1.05)


def price_rows(s, items, top):
    """Прайс: название слева, цена акцентом справа, между строками тонкая линия."""
    row_h = Inches(0.62)
    for i, (label, value) in enumerate(items):
        if i:
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PAD_X, top, CONTENT_W, Emu(9525))
            line.fill.solid()
            line.fill.fore_color.rgb = BORDER
            line.line.fill.background()
            line.shadow.inherit = False

        cell = s.shapes.add_textbox(PAD_X, top + Inches(0.09), CONTENT_W, row_h)
        tf = cell.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        run = tf.paragraphs[0].add_run()
        run.text = label
        run.font.size = pt(3.6)
        run.font.color.rgb = TEXT
        run.font.name = FONT

        num = s.shapes.add_textbox(PAD_X, top + Inches(0.05), CONTENT_W, row_h)
        tf = num.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(value)
        run.font.size = pt(4.6)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        run.font.name = FONT

        top += row_h
    return top


def card(s, top, height):
    """Подложка под выделенный блок — как .card в HTML."""
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PAD_X, top, CONTENT_W, height)
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.fill.background()
    box.shadow.inherit = False
    box.adjustments[0] = 0.06
    return box


def bullets(s, items, top):
    for item in items:
        dot = s.shapes.add_shape(
            MSO_SHAPE.OVAL, PAD_X, top + Inches(0.13), Inches(0.09), Inches(0.09)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        dot.shadow.inherit = False
        top = text(
            s, item, top, pt(3.4), left=PAD_X + Inches(0.28),
            width=CONTENT_W - Inches(0.28), spacing=1.35,
        ) + Inches(0.12)
    return top


# ─── 1. Титул ────────────────────────────────────────────────────────────────
s = slide()
text(s, "ЛАГЕРЬ KEEP", Inches(2.5), pt(2.6), color=MUTED)
text(s, "Искры", Inches(2.9), pt(16), bold=True, spacing=0.95)
text(s, "keep-sparks.ru", Inches(4.9), pt(4.2), color=MUTED)

# ─── 2. Что это ──────────────────────────────────────────────────────────────
s = slide()
title(s, "Что это")
text(s, "Искры — очки, которые ты набираешь за смену.", Inches(1.9), pt(4.2))
facts = [
    "Копятся через все смены\nза все годы",
    "Ничего не сгорает: приехал\nчерез год — продолжил\nс того же места",
    "По ним строится\nобщий рейтинг лагеря",
]
gap, card_w = Inches(0.3), (CONTENT_W - 2 * Inches(0.3)) / 3
for i, f in enumerate(facts):
    left = PAD_X + i * (card_w + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.2), card_w, Inches(2.1))
    box.fill.solid()
    box.fill.fore_color.rgb = SURFACE
    box.line.fill.background()
    box.shadow.inherit = False
    box.adjustments[0] = 0.08
    text(s, f, Inches(3.5), pt(3.2), left=left + Inches(0.3), width=card_w - Inches(0.6), spacing=1.35)

# ─── 3. День присутствия ─────────────────────────────────────────────────────
s = slide()
title(s, "Просто за то, что ты здесь")
text(s, "30", Inches(1.7), pt(26), color=ACCENT, bold=True, spacing=0.9)
text(s, "искр за каждый день смены, кроме дня отъезда.", Inches(5.0), pt(4.2))
text(s, "Всем и без условий. Делать для этого ничего не нужно.", Inches(5.6), pt(4.2), color=MUTED)

# ─── 4. Командные ────────────────────────────────────────────────────────────
s = slide()
title(s, "Всей командой")
bottom = price_rows(
    s,
    [
        ("КТБ: этап", 250),
        ("КТБ: победа", 400),
        ("КГГ/КТП: кубок", 150),
        ("КГГ/КТП: победа", 500),
        ("Wake Up Арена: победа комнаты", 300),
    ],
    Inches(1.75),
)
text(
    s,
    "Эти искры получает каждый в команде — не капитан\nи не только тот, кто выходил на этап.",
    bottom + Inches(0.35),
    pt(3.6),
    spacing=1.35,
)

# ─── 5. Личное в команде ─────────────────────────────────────────────────────
s = slide()
title(s, "Личное внутри команды")
bottom = price_rows(
    s, [("КТБ: лучший в команде", 600), ("КГГ/КТП: лучший из лучших", 1500)], Inches(2.1)
)
text(
    s,
    "Wake Up Арена играется комнатами по 5–6 человек,\n4 раунда за смену. Выиграла комната — искры каждому\nеё жителю.",
    bottom + Inches(0.5),
    pt(3.4),
    color=MUTED,
    spacing=1.35,
)

# ─── 6. Реалити ──────────────────────────────────────────────────────────────
s = slide()
title(s, "Реалити")
bottom = price_rows(
    s,
    [
        ("Победа", 3000),
        ("Супер-финал", 500),
        ("Финал", 200),
        ("Продвинул сюжет", 250),
        ("Лучший / лидер дня", 80),
    ],
    Inches(1.5),
)
card_top = bottom + Inches(0.3)
card(s, card_top, Inches(1.5))
box = s.shapes.add_textbox(PAD_X + Inches(0.35), card_top + Inches(0.25), CONTENT_W, Inches(0.5))
tf = box.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
for chunk, color, bold in [("3000 + 500 + 200 = ", TEXT, False), ("3700", ACCENT, True)]:
    run = p.add_run()
    run.text = chunk
    run.font.size = pt(5.4)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
text(
    s,
    "Победитель прошёл и финал, и супер-финал, поэтому получает всё сразу.",
    card_top + Inches(0.9),
    pt(3.0),
    color=MUTED,
    left=PAD_X + Inches(0.35),
    width=CONTENT_W - Inches(0.7),
)

# ─── 7. Личные награды ───────────────────────────────────────────────────────
s = slide()
title(s, "Личные награды")
price_rows(
    s,
    [
        ("Человек смены", 1300),
        ("Признание руководителя", 1200),
        ("Человек дня", 300),
        ("Звёзды: победа", 2000),
        ("Звёзды: финал", 600),
    ],
    Inches(1.9),
)

# ─── 8. Коэффициент ──────────────────────────────────────────────────────────
s = slide()
title(s, "Коэффициент смены")
text(
    s,
    "Чем больше народу приехало, тем труднее выделиться —\nи тем дороже смена. На этой смене 52 человека.",
    Inches(1.9),
    pt(4.0),
    spacing=1.35,
)
card(s, Inches(3.6), Inches(1.9))
box = s.shapes.add_textbox(PAD_X + Inches(0.35), Inches(3.9), CONTENT_W, Inches(0.6))
tf = box.text_frame
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
p = tf.paragraphs[0]
for chunk, color, bold in [
    ("1000 искр за смену → ", TEXT, False),
    ("1720", ACCENT, True),
    (" в профиле", TEXT, False),
]:
    run = p.add_run()
    run.text = chunk
    run.font.size = pt(5.4)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
text(
    s,
    "Умножается вся сумма за смену целиком, округление одно. Поэтому число\nв профиле может на единицу отличаться от сложенных вручную цен.",
    Inches(4.6),
    pt(3.0),
    color=MUTED,
    left=PAD_X + Inches(0.35),
    width=CONTENT_W - Inches(0.7),
    spacing=1.35,
)

# ─── 9. Раскрытие по дням ────────────────────────────────────────────────────
s = slide()
title(s, "Искры приходят по дням")
text(
    s,
    "Итоги дня подводят взрослые. Когда день отдан,\nна сайте появляется карточка.",
    Inches(1.9),
    pt(4.0),
    spacing=1.35,
)
bullets(
    s,
    [
        "Сначала — только «тебе пришли искры за вчера»",
        "Состав открывается по нажатию",
        "В общий счёт день идёт сразу, открыл ты карточку или нет",
    ],
    Inches(3.6),
)

# ─── 10. Где смотреть ────────────────────────────────────────────────────────
s = slide()
title(s, "Где смотреть")
bullets(
    s,
    [
        "Профиль — все искры и график по сменам",
        "Моя смена — искры по дням и своя команда КТБ",
        "Рейтинг искр — где ты среди всех",
        "За что искры — эти же цены, всегда под рукой",
        "Победители, человек дня, человек смены",
    ],
    Inches(2.0),
)

# ─── 11. QR ──────────────────────────────────────────────────────────────────
s = slide()
text(s, "Заходи", Inches(0.5), pt(7.5), bold=True, align=PP_ALIGN.CENTER, spacing=1.05)
qr = Inches(3.45)
s.shapes.add_picture("qr-keep-sparks.png", (W - qr) / 2, Inches(1.5), qr, qr)
text(s, "keep-sparks.ru", Inches(5.25), pt(4.4), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
text(s, "Логин и пароль — у вожатого", Inches(6.0), pt(3.4), color=MUTED, align=PP_ALIGN.CENTER)

# ─── 12. Финал ───────────────────────────────────────────────────────────────
s = slide()
text(s, "Считай свои искры", Inches(2.7), pt(7.5), bold=True, align=PP_ALIGN.CENTER, spacing=1.05)
text(s, "keep-sparks.ru", Inches(3.9), pt(6.0), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Слайды собираются сверху вниз, а показываются по центру — как в HTML-версии,
# где у секции `justify-content: center`. Одним проходом сдвигаем содержимое
# каждого слайда так, чтобы поля сверху и снизу сравнялись; взаимное
# расположение блоков при этом не меняется.
for s in prs.slides:
    shapes = [sh for sh in s.shapes]
    if not shapes:
        continue
    top = min(sh.top for sh in shapes)
    bottom = max(sh.top + sh.height for sh in shapes)
    shift = (H - (bottom - top)) // 2 - top
    for sh in shapes:
        sh.top += shift

prs.save("sparks.pptx")
print(f"sparks.pptx: {len(prs.slides._sldIdLst)} слайдов")
